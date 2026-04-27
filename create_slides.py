from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Palette de couleurs ─────────────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x1A, 0x37, 0x6C)   # titres
BLEU_CLAIR  = RGBColor(0x27, 0x6F, 0xBF)   # accents
VERT        = RGBColor(0x2E, 0x86, 0x5A)   # positif / résultats
ORANGE      = RGBColor(0xE0, 0x7B, 0x22)   # LSTM
VIOLET      = RGBColor(0x6A, 0x3D, 0x9A)   # Chronos
GRIS_CLAIR  = RGBColor(0xF2, 0xF4, 0xF8)   # fond boites
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
NOIR        = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # layout entièrement vide


# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=NOIR,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_bullet_box(slide, lines, x, y, w, h, bg=GRIS_CLAIR,
                   font_size=15, title=None, title_color=BLEU_FONCE):
    add_rect(slide, x, y, w, h, bg)
    tf_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12),
                                       Inches(w - 0.3), Inches(h - 0.2))
    tf_box.word_wrap = True
    tf = tf_box.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.bold = True
        r.font.size = Pt(font_size + 1)
        r.font.color.rgb = title_color
        r.font.name = "Calibri"
        first = False

    for line in lines:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.color.rgb = NOIR
        r.font.name = "Calibri"


def slide_background(slide, color=BLANC):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_image_safe(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        add_rect(slide, x, y, w, h, GRIS_CLAIR)
        add_text(slide, "[image non trouvée]", x + 0.1, y + h/2 - 0.2, w - 0.2, 0.4,
                 font_size=11, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
slide_background(slide, BLANC)

# Bande bleue en haut
add_rect(slide, 0, 0, 13.33, 4.5, BLEU_FONCE)

# Titre principal
add_text(slide,
         "Prévision du niveau des nappes\nphréatiques par réseaux de neurones",
         0.6, 0.6, 12.1, 2.5,
         font_size=36, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# Sous-titre
add_text(slide,
         "Données hydrologiques australiennes — 2 510 puits",
         0.6, 3.0, 12.1, 0.7,
         font_size=20, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)

# Ligne de séparation dorée
add_rect(slide, 2.5, 4.3, 8.3, 0.06, RGBColor(0xF0, 0xC0, 0x40))

# Infos bas de slide
add_text(slide, "Marius", 0.6, 4.7, 6, 0.6,
         font_size=16, bold=True, color=BLEU_FONCE)
add_text(slide, "Master 1 — 2025/2026", 6.8, 4.7, 6, 0.6,
         font_size=16, color=BLEU_FONCE, align=PP_ALIGN.RIGHT)

# Pastilles de parties
for i, (label, col) in enumerate([("Partie 1 : Visualisation", BLEU_CLAIR),
                                    ("Partie 2a : LSTM", ORANGE),
                                    ("Partie 2b : Chronos 2", VIOLET)]):
    x = 0.6 + i * 4.2
    add_rect(slide, x, 5.6, 3.8, 0.7, col)
    add_text(slide, label, x + 0.1, 5.65, 3.6, 0.6,
             font_size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Données & Visualisation
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
slide_background(slide, BLANC)

# Bandeau titre
add_rect(slide, 0, 0, 13.33, 1.1, BLEU_FONCE)
add_text(slide, "Partie 1 — Les données & Visualisation",
         0.3, 0.1, 12.7, 0.9,
         font_size=26, bold=True, color=BLANC, align=PP_ALIGN.LEFT)

# Colonne gauche — description données
add_bullet_box(slide,
    ["● 2 510 fichiers CSV (un par puit)",
     "● 6 variables par puit :",
     "     GWL  →  niveau de la nappe (m)  [cible]",
     "     P    →  précipitations (mm/mois)",
     "     T    →  température (°C)",
     "     ET   →  évapotranspiration (mm/mois)",
     "     NDVI →  indice de végétation",
     "     Date →  mensuelle",
     "● OUVRAGES.csv : coordonnées GPS + région",
     "● 8 régions : VIC, NSW, WA, QLD, SA, TAS, NT, ACT"],
    x=0.3, y=1.25, w=4.5, h=5.8,
    title="Les données", font_size=13)

# Colonne centre — méthode de visualisation
add_bullet_box(slide,
    ["● Fonction générique plot_well(id)",
     "",
     "● Figure A — Valeurs brutes",
     "  5 sous-graphiques superposés",
     "  (un par variable)",
     "",
     "● Figure B — Valeurs normalisées",
     "  Normalisation z-score :",
     "  x̃ = (x − μ) / σ",
     "  Permet de comparer des",
     "  variables d'unités différentes"],
    x=5.0, y=1.25, w=3.8, h=5.8,
    title="Méthode", font_size=13)

# Image exemple graphique
add_image_safe(slide, "graphiques/10029257_brut.png", 9.0, 1.25, 4.0, 2.7)
add_image_safe(slide, "graphiques/10029257_normalise.png", 9.0, 4.1, 4.0, 2.7)

add_text(slide, "Puit NSW — brut", 9.0, 3.9, 4.0, 0.3,
         font_size=9, color=BLEU_CLAIR, align=PP_ALIGN.CENTER)
add_text(slide, "Puit NSW — normalisé", 9.0, 6.75, 4.0, 0.3,
         font_size=9, color=BLEU_CLAIR, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LSTM
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
slide_background(slide, BLANC)

add_rect(slide, 0, 0, 13.33, 1.1, ORANGE)
add_text(slide, "Partie 2a — Modèle LSTM générique",
         0.3, 0.1, 12.7, 0.9,
         font_size=26, bold=True, color=BLANC)

# Architecture
add_bullet_box(slide,
    ["● Un seul modèle pour tous les 2 510 puits",
     "",
     "● Entrées (6 features) :",
     "  P, T, ET, NDVI + lat + lon",
     "",
     "● Séquences glissantes de 12 mois",
     "  → prédit GWL au mois suivant",
     "",
     "● Architecture finale (grid search) :",
     "  hidden = 128  |  layers = 3",
     "  lr = 5×10⁻⁴   |  batch = 512",
     "  epochs = 30   |  dropout = 0.2"],
    x=0.3, y=1.25, w=4.3, h=5.8,
    title="Architecture", font_size=13, title_color=ORANGE)

# Grid search
add_bullet_box(slide,
    ["● 108 combinaisons testées :",
     "  hidden : 64, 128, 256",
     "  layers : 1, 2, 3",
     "  lr     : 1e-3, 5e-4, 1e-4",
     "  batch  : 256, 512",
     "  epochs : 30, 50",
     "",
     "● Données préchargées en VRAM",
     "  (RTX 5070 Ti — 16 Go)",
     "  → batching sans transfert CPU"],
    x=4.8, y=1.25, w=4.0, h=5.8,
    title="Grid Search", font_size=13, title_color=ORANGE)

# Résultats LSTM
add_rect(slide, 9.1, 1.25, 3.9, 2.8, RGBColor(0xFF, 0xF0, 0xE0))
add_text(slide, "Résultats", 9.2, 1.3, 3.7, 0.45,
         font_size=15, bold=True, color=ORANGE)
for i, (label, val) in enumerate([
    ("RMSE moyen",  "0.805 m"),
    ("MAE moyen",   "0.686 m"),
    ("RMSE médian", "0.440 m"),
    ("MAE médian",  "0.356 m"),
]):
    add_text(slide, label, 9.2, 1.85 + i * 0.5, 2.2, 0.45, font_size=13, color=NOIR)
    add_text(slide, val,   11.3, 1.85 + i * 0.5, 1.6, 0.45,
             font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

# Schéma LSTM simplifié
add_rect(slide, 9.1, 4.3, 3.9, 2.6, GRIS_CLAIR)
add_text(slide, "Fonctionnement LSTM", 9.2, 4.35, 3.7, 0.4,
         font_size=13, bold=True, color=BLEU_FONCE)
add_text(slide,
         "12 mois d'historique\n"
         "[P, T, ET, NDVI, lat, lon]\n"
         "         ↓\n"
         "  3 couches LSTM\n"
         "         ↓\n"
         " Prédiction GWL (m)",
         9.2, 4.8, 3.7, 2.0,
         font_size=12, color=NOIR, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Chronos 2
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
slide_background(slide, BLANC)

add_rect(slide, 0, 0, 13.33, 1.1, VIOLET)
add_text(slide, "Partie 2b — Chronos 2 (Foundation Model)",
         0.3, 0.1, 12.7, 0.9,
         font_size=26, bold=True, color=BLANC)

# Qu'est-ce que Chronos 2
add_bullet_box(slide,
    ["● Modèle pré-entraîné par Amazon sur",
     "  des milliers de séries temporelles",
     "",
     "● Apprentissage auto-supervisé",
     "  (pas besoin de nos données)",
     "",
     "● Utilisation zero-shot :",
     "  aucun entraînement sur nos puits",
     "",
     "● Entrée : historique GWL uniquement",
     "  (modèle univarié)",
     "",
     "● Prédit une distribution probabiliste",
     "  → on prend la médiane"],
    x=0.3, y=1.25, w=4.3, h=5.8,
    title="Qu'est-ce que Chronos 2 ?", font_size=13, title_color=VIOLET)

# Résultats Chronos
add_rect(slide, 4.8, 1.25, 4.0, 2.8, RGBColor(0xF3, 0xEC, 0xFF))
add_text(slide, "Résultats — 1 174 puits", 4.9, 1.3, 3.8, 0.45,
         font_size=15, bold=True, color=VIOLET)
for i, (label, val) in enumerate([
    ("RMSE moyen",  "0.673 m"),
    ("MAE moyen",   "0.578 m"),
    ("RMSE médian", "0.301 m"),
    ("MAE médian",  "0.248 m"),
]):
    add_text(slide, label, 4.9, 1.85 + i * 0.5, 2.2, 0.45, font_size=13, color=NOIR)
    add_text(slide, val, 7.0, 1.85 + i * 0.5, 1.7, 0.45,
             font_size=13, bold=True, color=VIOLET, align=PP_ALIGN.RIGHT)

# Limites
add_bullet_box(slide,
    ["● Ne voit pas P, T, ET, NDVI",
     "  (univarié vs multivarié pour LSTM)",
     "",
     "● Comparison partielle :",
     "  1174 puits vs 854 pour le LSTM",
     "  (seuils de données différents)",
     "",
     "● Modèle généraliste, pas spécialisé",
     "  nappes phréatiques"],
    x=4.8, y=4.2, w=4.0, h=2.8,
    title="Limites de la comparaison", font_size=13, title_color=VIOLET)

# Fonctionnement
add_rect(slide, 9.1, 1.25, 3.9, 5.8, GRIS_CLAIR)
add_text(slide, "Fonctionnement", 9.2, 1.3, 3.7, 0.45,
         font_size=14, bold=True, color=VIOLET)
add_text(slide,
         "Historique GWL\n(train = tout - 12 mois)\n"
         "         ↓\n"
         "  chronos-bolt-small\n"
         "  (pré-entraîné)\n"
         "         ↓\n"
         "  N scénarios possibles\n"
         "  pour les 12 prochains mois\n"
         "         ↓\n"
         "  Médiane des scénarios\n"
         "  = prédiction finale",
         9.2, 1.85, 3.7, 4.8,
         font_size=12, color=NOIR, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Conclusion & Comparaison
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
slide_background(slide, BLANC)

add_rect(slide, 0, 0, 13.33, 1.1, VERT)
add_text(slide, "Conclusion — Comparaison des approches",
         0.3, 0.1, 12.7, 0.9,
         font_size=26, bold=True, color=BLANC)

# Tableau comparatif
headers = ["", "LSTM (2a)", "Chronos 2 (2b)", "Gain"]
rows = [
    ["RMSE moyen",  "0.805 m", "0.673 m", "−16 %"],
    ["MAE moyen",   "0.686 m", "0.578 m", "−16 %"],
    ["RMSE médian", "0.440 m", "0.301 m", "−32 %"],
    ["MAE médian",  "0.356 m", "0.248 m", "−30 %"],
]

col_x = [0.3, 3.1, 6.2, 9.8]
col_w = [2.7, 2.9, 3.4, 2.2]

# En-têtes
add_rect(slide, 0.3, 1.3, 11.7, 0.55, BLEU_FONCE)
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_text(slide, h, x + 0.1, 1.33, w, 0.48,
             font_size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# Lignes du tableau
for i, row in enumerate(rows):
    bg = GRIS_CLAIR if i % 2 == 0 else BLANC
    add_rect(slide, 0.3, 1.85 + i * 0.55, 11.7, 0.55, bg)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        color = VERT if j == 3 else (ORANGE if j == 1 else (VIOLET if j == 2 else NOIR))
        bold  = j in (1, 2, 3)
        add_text(slide, cell, x + 0.1, 1.88 + i * 0.55, w - 0.1, 0.48,
                 font_size=14, bold=bold, color=color, align=PP_ALIGN.CENTER)

# Conclusion LSTM
add_bullet_box(slide,
    ["● Modèle entraîné sur 2 510 puits",
     "● Utilise 6 features (météo + GPS)",
     "● Grid search : 108 configurations",
     "● Bon résultat pour un modèle from scratch"],
    x=0.3, y=4.15, w=5.6, h=2.9,
    title="LSTM — Points clés", font_size=13, title_color=ORANGE)

# Conclusion Chronos
add_bullet_box(slide,
    ["● Zero-shot : aucun entraînement",
     "● Plus précis malgré moins d'infos",
     "● Foundation models : très efficaces",
     "  sur les séries temporelles régulières"],
    x=6.1, y=4.15, w=5.6, h=2.9,
    title="Chronos 2 — Points clés", font_size=13, title_color=VIOLET)

# Message clé
add_rect(slide, 0.3, 7.0, 12.73, 0.38, RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide,
         "Les foundation models pré-entraînés peuvent surpasser des modèles spécialisés "
         "sans nécessiter d'entraînement sur les données cibles.",
         0.5, 7.02, 12.3, 0.35,
         font_size=12, bold=True, color=VERT, align=PP_ALIGN.CENTER)


# ─── Sauvegarde ──────────────────────────────────────────────────────────────
OUTPUT = "presentation_nappes.pptx"
prs.save(OUTPUT)
print(f"Présentation créée : {OUTPUT}")
